"""Gera apresentacao de kickoff OLICMAT para stakeholders."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Theme ──────────────────────────────────────────────
BG_DARK = RGBColor(0x0A, 0x0A, 0x14)       # fundo principal
BG_CARD = RGBColor(0x12, 0x12, 0x24)       # cartao / destaque
ACCENT_GOLD = RGBColor(0xD4, 0xA8, 0x3C)   # ouro OLICMAT
ACCENT_TEAL = RGBColor(0x2D, 0xD4, 0xBF)   # teal para detalhes
TEXT_WHITE = RGBColor(0xF0, 0xF0, 0xF0)
TEXT_MUTED = RGBColor(0xA0, 0xA0, 0xB0)
TEXT_BODY = RGBColor(0xCC, 0xCC, 0xDD)

W = Inches(13.333)  # 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── Helpers ────────────────────────────────────────────

def slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Outfit"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=TEXT_BODY, line_spacing=1.5, font_name="Outfit"):
    """lines = list of (text, bold, size_override, color_override)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, bold, size, clr = line, False, font_size, color
        else:
            text = line[0]
            bold = line[1] if len(line) > 1 else False
            size = line[2] if len(line) > 2 else font_size
            clr = line[3] if len(line) > 3 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(size * line_spacing * 0.4)
    return tf

def add_card(slide, left, top, width, height, color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    return shape

def add_line(slide, left, top, width, color=ACCENT_GOLD, thickness=Pt(2)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_icon_text(slide, left, top, icon, title, body, icon_color=ACCENT_GOLD):
    """Card com icone grande + titulo + descricao"""
    add_card(slide, left, top, 3.8, 2.8, border_color=ACCENT_GOLD)
    add_textbox(slide, left + 0.3, top + 0.2, 3.2, 0.6, icon,
                font_size=32, color=icon_color, bold=True)
    add_textbox(slide, left + 0.3, top + 0.8, 3.2, 0.5, title,
                font_size=18, color=TEXT_WHITE, bold=True)
    add_textbox(slide, left + 0.3, top + 1.4, 3.2, 1.2, body,
                font_size=13, color=TEXT_MUTED)

def slide_number(slide, num):
    add_textbox(slide, 12.2, 7.0, 0.8, 0.4, str(num),
                font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)

def add_bottom_bar(slide):
    add_line(slide, 0.8, 6.85, 11.7, color=ACCENT_GOLD, thickness=Pt(1))

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — Capa
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide_bg(sl)

# Gradiente simulado com cards
add_card(sl, 0, 0, 13.333, 7.5, color=BG_DARK)
add_line(sl, 1.5, 2.8, 3.0, color=ACCENT_GOLD, thickness=Pt(4))

add_textbox(sl, 1.5, 1.5, 10.5, 0.8, "OLICMAT 2026",
            font_size=56, color=TEXT_WHITE, bold=True)
add_textbox(sl, 1.5, 2.2, 10.5, 0.5,
            "Olimpíada para Licenciandos em Matemática",
            font_size=24, color=ACCENT_GOLD)
add_textbox(sl, 1.5, 3.2, 10.5, 1.0,
            "Plataforma integrada de competição olímpica, formação pedagógica\ne congresso acadêmico para futuros professores de Matemática",
            font_size=16, color=TEXT_MUTED)

add_textbox(sl, 1.5, 5.0, 5.0, 0.4, "Maio de 2026",
            font_size=14, color=TEXT_MUTED)
add_textbox(sl, 1.5, 5.4, 5.0, 0.4, "Kickoff — Equipe OLICMAT",
            font_size=14, color=TEXT_MUTED)

slide_number(sl, 1)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_textbox(sl, 1.0, 0.6, 8.0, 0.7, "Agenda", font_size=36, bold=True)
add_line(sl, 1.0, 1.2, 2.0, color=ACCENT_GOLD, thickness=Pt(3))

agenda = [
    ("1.", "O que é a OLICMAT", "Visão geral e os três pilares"),
    ("2.", "Público-alvo e impacto", "Quem participa e o que ganha"),
    ("3.", "Cronograma 2026", "Marcos e datas importantes"),
    ("4.", "A Plataforma", "O que está pronto e como funciona"),
    ("5.", "Funcionalidades", "Destaques por pilar"),
    ("6.", "Próximos passos", "Entregas, riscos e o que precisamos"),
]

for i, (num, title, desc) in enumerate(agenda):
    y = 2.0 + i * 0.85
    add_textbox(sl, 1.5, y, 0.5, 0.5, num, font_size=20, color=ACCENT_GOLD, bold=True)
    add_textbox(sl, 2.2, y, 8.0, 0.4, title, font_size=20, color=TEXT_WHITE, bold=True)
    add_textbox(sl, 2.2, y + 0.35, 8.0, 0.35, desc, font_size=13, color=TEXT_MUTED)

slide_number(sl, 2)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — Visão Geral
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_textbox(sl, 1.0, 0.6, 10.0, 0.7, "O que é a OLICMAT", font_size=36, bold=True)
add_line(sl, 1.0, 1.2, 2.0, color=ACCENT_GOLD, thickness=Pt(3))

add_textbox(sl, 1.0, 1.8, 11.0, 1.2,
            "Uma plataforma digital que une três iniciativas acadêmicas em uma só experiência, "
            "voltada exclusivamente para estudantes de Licenciatura em Matemática "
            "de instituições públicas brasileiras.",
            font_size=16, color=TEXT_BODY)

# 3 pilares como cards
add_icon_text(sl, 1.0, 3.5, "🏆", "OLICMAT — Competição",
              "Olimpíada em duas fases: prova objetiva online (30 questões) + "
              "produção de videoaula e portfólio digital. Ranking estadual com "
              "medalhas de ouro, prata e bronze.",
              icon_color=ACCENT_GOLD)

add_icon_text(sl, 5.2, 3.5, "📚", "FORPEMAT — Formação",
              "Programa de formação pedagógica com 14 módulos (120 horas). "
              "Conteúdo 100% online, progresso individual e certificado automático "
              "com código de validação único.",
              icon_color=ACCENT_TEAL)

add_icon_text(sl, 9.4, 3.5, "📖", "CONGEMAT — Congresso",
              "Congresso acadêmico para submissão de artigos completos e pôsteres. "
              "Avaliação por pares, fluxo digital do início ao fim, sem custos "
              "para os participantes.",
              icon_color=RGBColor(0xC0, 0x60, 0xA0))

add_bottom_bar(sl)
slide_number(sl, 3)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — Público-Alvo
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_textbox(sl, 1.0, 0.6, 10.0, 0.7, "Público-Alvo e Impacto", font_size=36, bold=True)
add_line(sl, 1.0, 1.2, 2.0, color=ACCENT_GOLD, thickness=Pt(3))

# Stats
stats = [
    ("~ 65.000", "Licenciandos em\nMatemática no Brasil", ACCENT_GOLD),
    ("186", "Instituições públicas\ncom o curso", ACCENT_TEAL),
    ("27", "Estados brasileiros\nabrangidos", RGBColor(0xC0, 0x60, 0xA0)),
    ("1.000+", "Inscritos esperados\nna 1ª edição", ACCENT_GOLD),
]

for i, (num, label, color) in enumerate(stats):
    x = 1.0 + i * 3.1
    add_card(sl, x, 2.0, 2.8, 2.0, border_color=color)
    add_textbox(sl, x + 0.2, 2.3, 2.4, 0.6, num, font_size=38, color=color, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(sl, x + 0.2, 3.0, 2.4, 0.8, label, font_size=13, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

# Benefícios
add_textbox(sl, 1.0, 4.5, 11.0, 0.5, "O que os participantes ganham:",
            font_size=18, color=TEXT_WHITE, bold=True)

benefits = [
    ("Medalhas e premiação", "Reconhecimento acadêmico com ranking estadual"),
    ("Certificado 120 horas", "Formação pedagógica complementar gratuita"),
    ("Publicação acadêmica", "Submissão de artigos sem custo"),
    ("Currículo fortalecido", "Diferencial para mestrado e seleções"),
]

for i, (title, desc) in enumerate(benefits):
    x = 1.0 + (i % 2) * 6.1
    y = 5.1 + (i // 2) * 0.75
    add_textbox(sl, x, y, 0.4, 0.4, "▸", font_size=14, color=ACCENT_TEAL)
    add_textbox(sl, x + 0.4, y, 5.5, 0.35, title, font_size=15, color=TEXT_WHITE, bold=True)
    add_textbox(sl, x + 0.4, y + 0.3, 5.5, 0.3, desc, font_size=12, color=TEXT_MUTED)

add_bottom_bar(sl)
slide_number(sl, 4)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — Cronograma 2026
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_textbox(sl, 1.0, 0.6, 10.0, 0.7, "Cronograma 2026", font_size=36, bold=True)
add_line(sl, 1.0, 1.2, 2.0, color=ACCENT_GOLD, thickness=Pt(3))

timeline = [
    ("MAI", "Plataforma no ar\nVersão beta", ACCENT_TEAL),
    ("JUN", "Abertura das\ninscrições", ACCENT_GOLD),
    ("JUL", "Fase 1\nProva online", ACCENT_GOLD),
    ("AGO", "Resultado F1 +\nInício Fase 2", ACCENT_TEAL),
    ("SET", "Envio da\nFase 2", ACCENT_GOLD),
    ("OUT", "Avaliação F2 +\nCONGEMAT", ACCENT_TEAL),
    ("NOV", "Resultado final\nMedalhas + Certificados", ACCENT_GOLD),
]

for i, (mes, desc, color) in enumerate(timeline):
    x = 0.7 + i * 1.75
    # circle
    circ = sl.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x + 0.4), Inches(2.5), Inches(0.7), Inches(0.7))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    add_textbox(sl, x + 0.05, 2.55, 1.4, 0.5, mes,
                font_size=16, color=BG_DARK if color == ACCENT_TEAL else BG_DARK,
                bold=True, alignment=PP_ALIGN.CENTER)
    # conector
    if i < len(timeline) - 1:
        add_line(sl, x + 1.15, 2.8, 1.3, color=ACCENT_GOLD, thickness=Pt(2))
    # desc
    add_textbox(sl, x - 0.1, 3.5, 1.8, 0.8, desc,
                font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Destaque: maio
add_card(sl, 1.0, 5.0, 11.3, 1.4, border_color=ACCENT_TEAL)
add_textbox(sl, 1.5, 5.2, 10.3, 0.4, "🟢 Momento atual: Maio 2026 — Plataforma em versão beta",
            font_size=16, color=ACCENT_TEAL, bold=True)
add_textbox(sl, 1.5, 5.7, 10.3, 0.5,
            "Sistema já contempla: autenticação completa, prova online, ranking, "
            "catálogo FORPEMAT, submissão CONGEMAT e upload de arquivos via Cloudinary.",
            font_size=13, color=TEXT_BODY)

add_bottom_bar(sl)
slide_number(sl, 5)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — A Plataforma (Arquitetura Simplificada)
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_textbox(sl, 1.0, 0.6, 10.0, 0.7, "A Plataforma", font_size=36, bold=True)
add_line(sl, 1.0, 1.2, 2.0, color=ACCENT_GOLD, thickness=Pt(3))

# Três camadas visuais
layers = [
    ("Frontend", "Next.js 16 + Tailwind CSS + shadcn/ui",
     "Interface moderna, responsiva, tema escuro, acessível em qualquer dispositivo",
     "🌐", ACCENT_TEAL),
    ("Backend", "NestJS 11 + Prisma + PostgreSQL 16",
     "API REST com autenticação JWT, validação Zod, upload Cloudinary",
     "⚙️", ACCENT_GOLD),
    ("Infra", "Docker + Cloudinary",
     "Containerizado, deploy simplificado, armazenamento de mídia na nuvem",
     "☁️", RGBColor(0xC0, 0x60, 0xA0)),
]

for i, (title, tech, desc, icon, color) in enumerate(layers):
    y = 2.0 + i * 1.6
    add_card(sl, 1.0, y, 11.3, 1.3, border_color=color)
    add_textbox(sl, 1.3, y + 0.15, 0.7, 0.5, icon, font_size=28, color=color)
    add_textbox(sl, 2.2, y + 0.15, 3.0, 0.4, title, font_size=20, color=TEXT_WHITE, bold=True)
    add_textbox(sl, 5.5, y + 0.15, 5.5, 0.4, tech, font_size=13, color=TEXT_MUTED)
    add_textbox(sl, 2.2, y + 0.7, 9.5, 0.4, desc, font_size=14, color=TEXT_BODY)

add_bottom_bar(sl)
slide_number(sl, 6)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — Destaques: OLICMAT (Competição)
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_textbox(sl, 1.0, 0.6, 10.0, 0.7, "Destaque: Competição OLICMAT", font_size=36, bold=True)
add_line(sl, 1.0, 1.2, 2.0, color=ACCENT_GOLD, thickness=Pt(3))

add_textbox(sl, 1.0, 1.7, 11.0, 0.5,
            "Duas fases que avaliam conhecimento teórico e habilidade didática:",
            font_size=15, color=TEXT_MUTED)

# Fase 1
add_card(sl, 1.0, 2.5, 5.5, 3.2, border_color=ACCENT_GOLD)
add_textbox(sl, 1.3, 2.7, 5.0, 0.4, "FASE 1 — Prova Objetiva", font_size=20, color=ACCENT_GOLD, bold=True)
f1_items = [
    "30 questões de múltipla escolha (A-E)",
    "Eixos: Álgebra, Geometria, Análise, Estatística, Didática",
    "Cronômetro regressivo de 180 minutos",
    "Correção automática ao finalizar",
    "Salvamento contínuo das respostas",
    "Peso: 40% da nota final",
]
for i, item in enumerate(f1_items):
    add_textbox(sl, 1.5, 3.3 + i * 0.38, 4.8, 0.35, f"• {item}", font_size=13, color=TEXT_BODY)

# Fase 2
add_card(sl, 7.0, 2.5, 5.5, 3.2, border_color=ACCENT_TEAL)
add_textbox(sl, 7.3, 2.7, 5.0, 0.4, "FASE 2 — Produção Didática", font_size=20, color=ACCENT_TEAL, bold=True)
f2_items = [
    "Sorteio de tema entre 10 opções",
    "Produção de videoaula (até 20 min, MP4)",
    "Elaboração de portfólio digital (PDF)",
    "Upload direto pela plataforma (Cloudinary)",
    "Avaliação por banca de professores",
    "Peso: 60% da nota final",
]
for i, item in enumerate(f2_items):
    add_textbox(sl, 7.5, 3.3 + i * 0.38, 4.8, 0.35, f"• {item}", font_size=13, color=TEXT_BODY)

# Resultado
add_card(sl, 1.0, 6.0, 11.3, 0.8, border_color=RGBColor(0xC0, 0x60, 0xA0))
add_textbox(sl, 1.5, 6.15, 10.3, 0.4,
            "🏅 Ranking estadual com medalhas: Ouro (top 5%), Prata (10%), Bronze (15%) — "
            "cálculo automático por estado",
            font_size=14, color=TEXT_WHITE, bold=True)

add_bottom_bar(sl)
slide_number(sl, 7)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 — Destaques: FORPEMAT + CONGEMAT
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_textbox(sl, 1.0, 0.6, 10.0, 0.7, "Destaque: Formação e Congresso", font_size=36, bold=True)
add_line(sl, 1.0, 1.2, 2.0, color=ACCENT_GOLD, thickness=Pt(3))

# FORPEMAT
add_card(sl, 1.0, 1.8, 5.5, 4.5, border_color=ACCENT_TEAL)
add_textbox(sl, 1.3, 2.0, 5.0, 0.4, "📚 FORPEMAT — 14 Módulos", font_size=18, color=ACCENT_TEAL, bold=True)
forpemat_mods = [
    "TDICs na Educação Matemática",
    "Metodologias Ativas no Ensino",
    "Resolução de Problemas",
    "Modelagem Matemática",
    "Didática da Matemática",
    "Avaliação da Aprendizagem",
    "Etnomatemática",
    "História da Matemática",
    "Geometria Dinâmica",
    "Pensamento Computacional",
    "Educação Inclusiva",
    "Currículo e BNCC",
    "Jogos no Ensino de Matemática",
    "Projeto Integrador: Sequência Didática Digital",
]
for i, mod in enumerate(forpemat_mods):
    add_textbox(sl, 1.5, 2.6 + i * 0.27, 4.8, 0.25, f"• {mod}", font_size=11, color=TEXT_BODY)

add_card(sl, 1.5, 6.2, 4.5, 0.5, color=ACCENT_TEAL)
add_textbox(sl, 1.7, 6.3, 4.3, 0.3,
            "Certificado automático de 120h ao concluir todos",
            font_size=12, color=BG_DARK, bold=True)

# CONGEMAT
add_card(sl, 7.0, 1.8, 5.5, 3.0, border_color=RGBColor(0xC0, 0x60, 0xA0))
add_textbox(sl, 7.3, 2.0, 5.0, 0.4, "📖 CONGEMAT — Congresso", font_size=18,
            color=RGBColor(0xC0, 0x60, 0xA0), bold=True)
congemat_items = [
    "Submissão de artigos completos e pôsteres",
    "Upload de PDF diretamente na plataforma",
    "Fluxo de avaliação por pares",
    "Status em tempo real: Em avaliação / Aprovado / Rejeitado",
    "Sem custos para participantes",
    "Anais digitais com todos os trabalhos aprovados",
]
for i, item in enumerate(congemat_items):
    add_textbox(sl, 7.5, 2.6 + i * 0.35, 4.8, 0.3, f"• {item}", font_size=12, color=TEXT_BODY)

# Métricas lado direito
add_card(sl, 7.0, 5.2, 5.5, 1.6, border_color=ACCENT_GOLD)
add_textbox(sl, 7.3, 5.4, 5.0, 0.4, "🎯 Metas 2026", font_size=18, color=ACCENT_GOLD, bold=True)
metrics = [
    "Meta de 1.000+ inscritos na 1ª edição",
    "70% concluírem ao menos 1 módulo FORPEMAT",
    "100+ submissões no CONGEMAT",
]
for i, m in enumerate(metrics):
    add_textbox(sl, 7.5, 5.9 + i * 0.35, 4.8, 0.3, f"▸ {m}", font_size=12, color=TEXT_BODY)

add_bottom_bar(sl)
slide_number(sl, 8)

# ═══════════════════════════════════════════════════════════
# SLIDE 9 — O que já está pronto
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_textbox(sl, 1.0, 0.6, 10.0, 0.7, "Status Atual do Desenvolvimento", font_size=36, bold=True)
add_line(sl, 1.0, 1.2, 2.0, color=ACCENT_GOLD, thickness=Pt(3))

add_textbox(sl, 1.0, 1.6, 11.0, 0.4,
            "A plataforma está funcional — o MVP cobre os fluxos principais dos três pilares:",
            font_size=15, color=TEXT_MUTED)

done = [
    ("✅", "Autenticação completa", "Registro em 2 etapas, login JWT, 3 perfis (Aluno, Avaliador, Admin)"),
    ("✅", "Inscrição OLICMAT", "Formulário com UF, instituição, curso, período + validação por admin"),
    ("✅", "Prova Fase 1", "30 questões com timer de 180 min, salvamento contínuo, correção automática"),
    ("✅", "Upload Fase 2", "Videoaula (MP4) + Portfólio (PDF) via Cloudinary com barra de progresso"),
    ("✅", "Ranking + Medalhas", "Por estado, distribuição automática Ouro/Prata/Bronze (5%/10%/15%)"),
    ("✅", "FORPEMAT", "14 módulos com conteúdo, progresso individual, certificado automático 120h"),
    ("✅", "CONGEMAT", "Submissão de artigos/pôsteres, avaliação por pares, status em tempo real"),
    ("✅", "Docker", "Ambiente containerizado: PostgreSQL + Backend + Frontend"),
]

for i, (icon, title, desc) in enumerate(done):
    y = 2.3 + i * 0.65
    add_card(sl, 1.0, y, 11.3, 0.55, color=BG_CARD if i % 2 == 0 else RGBColor(0x18, 0x18, 0x30))
    add_textbox(sl, 1.2, y + 0.1, 0.4, 0.35, icon, font_size=16)
    add_textbox(sl, 1.7, y + 0.1, 4.0, 0.35, title, font_size=14, color=TEXT_WHITE, bold=True)
    add_textbox(sl, 5.5, y + 0.1, 6.5, 0.35, desc, font_size=12, color=TEXT_MUTED)

add_bottom_bar(sl)
slide_number(sl, 9)

# ═══════════════════════════════════════════════════════════
# SLIDE 10 — Próximos Passos
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_textbox(sl, 1.0, 0.6, 10.0, 0.7, "Próximos Passos", font_size=36, bold=True)
add_line(sl, 1.0, 1.2, 2.0, color=ACCENT_GOLD, thickness=Pt(3))

next_steps = [
    ("MAI", "Semanas 1-2", "Testes de carga e segurança", "Garantir que a plataforma suporte picos de acesso durante a prova"),
    ("MAI", "Semanas 3-4", "Ajustes finos e correções", "Revisão de UX, validação de formulários, mensagens de erro"),
    ("JUN", "Semana 1", "Lançamento oficial", "Abertura das inscrições e campanha de divulgação"),
    ("JUN", "Contínuo", "Suporte operacional", "Acompanhamento de inscrições, validação de documentos, helpdesk"),
]

for i, (mes, periodo, titulo, desc) in enumerate(next_steps):
    y = 1.8 + i * 1.25
    add_card(sl, 1.0, y, 11.3, 1.05, border_color=ACCENT_GOLD if i < 2 else ACCENT_TEAL)
    add_card(sl, 1.2, y + 0.15, 1.6, 0.7, color=ACCENT_GOLD if i < 2 else ACCENT_TEAL)
    add_textbox(sl, 1.3, y + 0.25, 1.4, 0.3, mes, font_size=16, color=BG_DARK, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(sl, 1.3, y + 0.55, 1.4, 0.25, periodo, font_size=10, color=BG_DARK,
                alignment=PP_ALIGN.CENTER)
    add_textbox(sl, 3.2, y + 0.15, 8.5, 0.35, titulo, font_size=18, color=TEXT_WHITE, bold=True)
    add_textbox(sl, 3.2, y + 0.55, 8.5, 0.35, desc, font_size=13, color=TEXT_MUTED)

# O que precisamos
add_card(sl, 1.0, 6.2, 11.3, 0.7, border_color=ACCENT_GOLD)
add_textbox(sl, 1.5, 6.3, 10.3, 0.4,
            "🔑 Para o lançamento: servidor de produção, domínio, contas oficiais Cloudinary, "
            "base de questões para prova e avaliadores confirmados",
            font_size=13, color=TEXT_WHITE)

add_bottom_bar(sl)
slide_number(sl, 10)

# ═══════════════════════════════════════════════════════════
# SLIDE 11 — Riscos
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_textbox(sl, 1.0, 0.6, 10.0, 0.7, "Riscos e Mitigações", font_size=36, bold=True)
add_line(sl, 1.0, 1.2, 2.0, color=ACCENT_GOLD, thickness=Pt(3))

risks = [
    ("⚠️", "Sobrecarga no dia da prova",
     "500+ alunos simultâneos podem sobrecarregar o servidor",
     "Testes de carga com k6, escalonamento horizontal, CDN para assets estáticos"),
    ("⚠️", "Queda de conexão durante a prova",
     "Aluno perder conexão e respostas não salvas",
     "Salvamento incremental via upsert a cada resposta, possibilidade de retomar"),
    ("⚠️", "Baixa adesão na primeira edição",
     "Menos de 500 inscritos",
     "Divulgação via instituições parceiras, redes sociais, coordenações de curso"),
    ("⚠️", "Atraso na disponibilidade dos avaliadores",
     "Fase 2 sem correção a tempo",
     "Recrutamento antecipado, plataforma de avaliação simplificada"),
]

for i, (icon, title, risk, mitigation) in enumerate(risks):
    y = 1.8 + i * 1.3
    add_card(sl, 1.0, y, 11.3, 1.1, border_color=RGBColor(0x80, 0x40, 0x40))
    add_textbox(sl, 1.3, y + 0.1, 0.4, 0.3, icon, font_size=18)
    add_textbox(sl, 1.8, y + 0.1, 4.0, 0.35, title, font_size=16, color=TEXT_WHITE, bold=True)
    add_textbox(sl, 1.8, y + 0.45, 4.0, 0.3, f"Risco: {risk}", font_size=11, color=RGBColor(0xE0, 0x80, 0x80))
    add_textbox(sl, 6.2, y + 0.15, 5.8, 0.7, f"Mitigação: {mitigation}", font_size=12, color=TEXT_BODY)

add_bottom_bar(sl)
slide_number(sl, 11)

# ═══════════════════════════════════════════════════════════
# SLIDE 12 — Encerramento
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_textbox(sl, 1.0, 1.5, 11.3, 0.7, "Vamos construir juntos a maior\nolimpíada de Matemática para licenciandos do Brasil.",
            font_size=36, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_line(sl, 4.5, 2.8, 4.3, color=ACCENT_GOLD, thickness=Pt(3))

add_textbox(sl, 2.0, 3.5, 9.3, 1.5,
            "A plataforma está pronta para entrar em produção.\n"
            "O cronograma é ambicioso mas factível.\n"
            "O impacto na formação de professores de Matemática será real e mensurável.",
            font_size=18, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)

add_textbox(sl, 2.0, 5.5, 9.3, 0.5, "Perguntas?",
            font_size=28, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(sl, 1.0, 6.5, 11.3, 0.4, "olicmat • contato@olicmat.com.br",
            font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

slide_number(sl, 12)

# ── Save ────────────────────────────────────────────────
output_path = "/home/wesley/Projetos/sites/olicmat/docs/Apresentacao_Kickoff_OLICMAT_2026.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"Apresentação salva em: {output_path}")
print(f"Total de slides: {len(prs.slides)}")
